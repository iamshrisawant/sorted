import hashlib
import logging
from pathlib import Path
from typing import Dict, Union

import docx
import numpy as np
import pandas as pd
import pdfplumber
from openpyxl import load_workbook
from pptx import Presentation
from sentence_transformers import SentenceTransformer

logging.getLogger("pdfminer").setLevel(logging.ERROR)
logger = logging.getLogger(__name__)

MODEL_NAME = "all-MiniLM-L6-v2"
_model = None
embedding_dim = 384

def _extract_content(path: Path, file_type: str) -> str:
    try:
        if file_type in ("pdf", "jpg", "jpeg", "png", "bmp", "webp"):
            from src.core.utils.extractor import extract_visual_content
            return extract_visual_content(path, file_type)
        elif file_type == "docx":
            doc = docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        elif file_type == "pptx":
            prs = Presentation(path)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif file_type == "xlsx":
            wb = load_workbook(path, data_only=True)
            return "\n".join(str(cell.value) for ws in wb.worksheets for row in ws.iter_rows() if cell.value is not None)
        elif file_type == "csv":
            df = pd.read_csv(path)
            return df.astype(str).to_string(index=False)
        elif file_type in ("tsv"):
            df = pd.read_csv(path, sep='\t')
            return df.astype(str).to_string(index=False)
        elif file_type in ("txt", "md", "mdx", "rst", "rtf", "html", "htm", "xml", "eml", "log"):
            return path.read_text(encoding="utf-8", errors="ignore")
        else:
            return ""
    except Exception as e:
        logger.error(f"[Processor] Failed to read {path.name} ({file_type}): {e}")
        return ""

def _compute_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()

def _load_model(local_only: bool = True):
    global _model
    if _model == "FAILED":
        raise RuntimeError("Model loading previously failed due to missing files or network errors. Skipping to prevent infinite crash loops.")
    
    if _model is None:
        try:
            from src.core.utils.paths import get_models_dir
            logger.info(f"[Processor] Loading model for the first time: {MODEL_NAME}")
            # Try to load local first. If not found, internet download fails gracefully.
            _model = SentenceTransformer(
                MODEL_NAME, 
                cache_folder=str(get_models_dir()), 
                local_files_only=local_only
            )
            logger.info("[Processor] Model loaded successfully.")
        except Exception as e:
            _model = "FAILED"
            logger.error(f"[Processor] Failed to load model: {e}")
            raise e
    return _model

def process_file(file_path: Union[str, Path]) -> Dict:
    path = Path(file_path)
    if not path.is_file():
        return {}
        
    file_type = path.suffix.lower().lstrip('.')
    raw_content = _extract_content(path, file_type)
    
    # Semantic Noise Filter: If visual/OCR extraction yields no text, abort.
    # This prevents indexing random visual junk (screenshots, OS icons) without semantics.
    if file_type in ("pdf", "jpg", "jpeg", "png", "bmp", "webp") and not raw_content.strip():
        logger.warning(f"[Processor] Skipping {path.name}: No semantic text extracted from visual data.")
        return {}

    text = raw_content.strip()
    if len(text) > 1000:
        cleaned_content = text[:500] + "\n...\n" + text[-500:]
    else:
        cleaned_content = text
        
    full_content = f"{path.name}\n{cleaned_content}"
    
    model = _load_model()
    emb = model.encode(full_content, show_progress_bar=False, normalize_embeddings=True)
    
    if isinstance(emb, np.ndarray):
        embeddings = [emb.tolist()]
    else:
        embeddings = []

    return {
        "file_path": str(path),
        "file_name": path.stem,
        "parent_folder": path.parent.name,
        "parent_folder_path": str(path.parent.resolve()),
        "file_type": file_type,
        "content_hash": _compute_hash(raw_content),
        "embeddings": embeddings,
    }

def process_text_context(text: str, target_folder: Union[str, Path]) -> Dict:
    path = Path(target_folder)
    if not path.exists() or not path.is_dir():
        logger.warning(f"[Processor] Target folder for context does not exist: {target_folder}")
        return {}
        
    cleaned_content = text.strip()
    full_content = f"{path.name} Context\n{cleaned_content}"
    
    model = _load_model()
    emb = model.encode(full_content, show_progress_bar=False, normalize_embeddings=True)
    
    if isinstance(emb, np.ndarray):
        embeddings = [emb.tolist()]
    else:
        embeddings = []

    # Create dummy file path representing the context
    dummy_path = path / "_folder_context_.txt"

    return {
        "file_path": str(dummy_path),
        "file_name": "_folder_context_",
        "parent_folder": path.name,
        "parent_folder_path": str(path.resolve()),
        "file_type": "txt",
        "content_hash": _compute_hash(cleaned_content),
        "embeddings": embeddings,
    }

